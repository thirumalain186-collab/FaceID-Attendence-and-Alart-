from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(10, 36, 99)      # #0A2463
PURPLE = RGBColor(36, 123, 160)       # #247BA0
CYAN = RGBColor(6, 214, 160)          # #06D6A0
RED = RGBColor(230, 57, 70)           # #E63946
GREEN = RGBColor(42, 157, 143)        # #2A9D8F
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
BLACK = RGBColor(0, 0, 0)

def add_background(slide, color):
    """Add solid background color to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide():
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, DARK_BLUE)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "AI-Powered Real-Time Attendance & Surveillance System"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "GPU Accelerated Multi-Person Tracking"
    p.font.size = Pt(32)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Built using
    built_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.5))
    built_frame = built_box.text_frame
    p = built_frame.paragraphs[0]
    p.text = "Built using Python + Electron"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Team names
    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1))
    team_frame = team_box.text_frame
    team_frame.word_wrap = True
    p = team_frame.paragraphs[0]
    p.text = "Team: Saravanan, Praveen, Harish Raj, Tharun Kumar, Sathivel"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Department and College
    dept_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1.2))
    dept_frame = dept_box.text_frame
    dept_frame.word_wrap = True
    p = dept_frame.paragraphs[0]
    p.text = "Department: B.Tech AI & Data Science\nCollege: P.T. Lee Chengalvaraya Naicker College of Engineering & Technology"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_problem_solution_slide():
    """Slide 2: Problem & Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, BLACK)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Problem & Solution"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Left side - Problem
    problem_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    problem_shape.fill.solid()
    problem_shape.fill.fore_color.rgb = RGBColor(60, 20, 20)  # Dark red
    problem_shape.line.color.rgb = RED
    problem_shape.line.width = Pt(2)
    
    problem_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(4.1), Inches(5.4))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True
    
    # Problem title
    p = problem_frame.paragraphs[0]
    p.text = "❌ PROBLEM"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RED
    
    # Problem points
    problems = [
        "Takes 5-10 mins per class",
        "20-30% proxy attendance",
        "No real-time alerts",
        "Cannot track 10+ people",
        "Human errors"
    ]
    
    for problem in problems:
        p = problem_frame.add_paragraph()
        p.text = "• " + problem
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)
    
    # Right side - Solution
    solution_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    solution_shape.fill.solid()
    solution_shape.fill.fore_color.rgb = RGBColor(20, 60, 20)  # Dark green
    solution_shape.line.color.rgb = GREEN
    solution_shape.line.width = Pt(2)
    
    solution_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.4), Inches(3.9), Inches(5.4))
    solution_frame = solution_box.text_frame
    solution_frame.word_wrap = True
    
    # Solution title
    p = solution_frame.paragraphs[0]
    p.text = "✅ SOLUTION"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    # Solution points
    solutions = [
        "Real-time processing",
        "AI face recognition",
        "Unknown person alerts",
        "Detects 20+ people",
        "Automated system"
    ]
    
    for solution in solutions:
        p = solution_frame.add_paragraph()
        p.text = "• " + solution
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(8)

def add_architecture_slide():
    """Slide 3: Architecture & Tech Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, BLACK)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Architecture & Tech Stack"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Left side - Architecture flow
    arch_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.8), Inches(5.8))
    arch_frame = arch_box.text_frame
    arch_frame.word_wrap = True
    
    # Add architecture layers
    layers = [
        ("💻 Electron UI", "Desktop App"),
        ("🔗 Flask API", "REST Endpoints"),
        ("🧠 Python AI Engine", "YOLO + Face Rec"),
        ("🗄️  Database", "SQLite + Alerts")
    ]
    
    y_pos = 0
    for layer_title, layer_desc in layers:
        p = arch_frame.paragraphs[0] if y_pos == 0 else arch_frame.add_paragraph()
        p.text = layer_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
        p = arch_frame.add_paragraph()
        p.text = layer_desc
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.level = 1
        
        if y_pos < len(layers) - 1:
            p = arch_frame.add_paragraph()
            p.text = "↓"
            p.font.size = Pt(14)
            p.font.color.rgb = CYAN
            p.alignment = PP_ALIGN.CENTER
        
        y_pos += 1
    
    # Right side - Tech Stack table
    tech_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    tech_frame = tech_box.text_frame
    tech_frame.word_wrap = True
    
    p = tech_frame.paragraphs[0]
    p.text = "🔧 TECH STACK"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(12)
    
    tech_stack = [
        ("Python 🐍", "AI Engine"),
        ("YOLO v8", "GPU Detection"),
        ("LBPH", "Face Recognition"),
        ("Electron ⚡", "Desktop UI"),
        ("Flask", "API Server"),
        ("SQLite", "Database"),
        ("GPU (CUDA)", "Acceleration")
    ]
    
    for tech, purpose in tech_stack:
        p = tech_frame.add_paragraph()
        p.text = f"{tech}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        p = tech_frame.add_paragraph()
        p.text = purpose
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.level = 1
        p.space_after = Pt(6)

def add_features_slide():
    """Slide 4: Features & Performance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, BLACK)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Key Features & Performance"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Left side - Features
    features_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.8), Inches(5.8))
    features_frame = features_box.text_frame
    features_frame.word_wrap = True
    
    p = features_frame.paragraphs[0]
    p.text = "🎯 FEATURES"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(12)
    
    features = [
        ("✅ 20+ People Detection", "Simultaneous"),
        ("⚡ 150ms per frame", "GPU Accelerated"),
        ("🎯 95% Accuracy", "Pre-trained model"),
        ("⚠️  Unknown Alerts", "Email + Photo"),
        ("💻 Easy Desktop UI", "One-click start"),
        ("🔒 Secure Database", "SQLite")
    ]
    
    for feature, detail in features:
        p = features_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        p = features_frame.add_paragraph()
        p.text = detail
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.level = 1
        p.space_after = Pt(8)
    
    # Right side - Performance Metrics
    perf_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    perf_frame = perf_box.text_frame
    perf_frame.word_wrap = True
    
    p = perf_frame.paragraphs[0]
    p.text = "📊 PERFORMANCE"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(12)
    
    metrics = [
        ("Our System", "Old System", "Manual"),
        ("95% Accuracy", "80% Accuracy", "70% Accuracy"),
        ("<30 seconds", "2-3 minutes", "5-10 minutes"),
        ("20+ people", "5-10 people", "❌"),
        ("Real-time alerts", "Delayed", "❌"),
        ("GPU Powered", "❌", "❌")
    ]
    
    for i, (ours, old, manual) in enumerate(metrics):
        if i == 0:
            p = perf_frame.paragraphs[0]
        else:
            p = perf_frame.add_paragraph()
        
        p.text = ours
        p.font.size = Pt(13) if i == 0 else Pt(12)
        p.font.bold = True if i == 0 else False
        p.font.color.rgb = CYAN if i == 0 else GREEN
        p.space_after = Pt(4)
        
        if i > 0:
            p = perf_frame.add_paragraph()
            p.text = f"vs {old} / {manual}"
            p.font.size = Pt(10)
            p.font.color.rgb = LIGHT_GRAY
            p.level = 1
            p.space_after = Pt(8)

def add_demo_slide():
    """Slide 5: Live Demo & Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, BLACK)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Live Demo & Results"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Demo sections
    demo_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    demo_frame = demo_box.text_frame
    demo_frame.word_wrap = True
    
    sections = [
        ("🎬 DEMONSTRATION", [
            "1. Electron Desktop UI launches",
            "2. Flask backend starts automatically",
            "3. Real-time face detection begins",
            "4. Unknown face triggers email alert",
            "5. Attendance marked automatically"
        ]),
        ("📊 REAL-TIME METRICS", [
            "• Detection Speed: 150ms per frame ⚡",
            "• Accuracy: 95%",
            "• Simultaneous People: 20+",
            "• Alert Delivery: <5 seconds",
            "• Database Query: <100ms"
        ])
    ]
    
    for section_title, items in sections:
        p = demo_frame.paragraphs[0] if demo_frame.paragraphs[0].text == "" else demo_frame.add_paragraph()
        p.text = section_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.space_after = Pt(8)
        
        for item in items:
            p = demo_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)
        
        p = demo_frame.add_paragraph()
        p.text = ""
        p.space_after = Pt(12)

def add_innovation_slide():
    """Slide 6: Innovation & Applications"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, BLACK)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Innovation & Applications"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Left side - Innovation
    innov_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.8), Inches(5.8))
    innov_frame = innov_box.text_frame
    innov_frame.word_wrap = True
    
    p = innov_frame.paragraphs[0]
    p.text = "🔥 INNOVATION"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(10)
    
    innovations = [
        ("Tracking + Recognition", "60% speed boost"),
        ("Auto Email Alerts", "<5 sec delivery"),
        ("GPU Optimization", "20+ people real-time"),
        ("Production Ready", "Deploy instantly")
    ]
    
    for innov, detail in innovations:
        p = innov_frame.add_paragraph()
        p.text = innov
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
        p = innov_frame.add_paragraph()
        p.text = detail
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.level = 1
        p.space_after = Pt(10)
    
    # Right side - Applications
    app_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    app_frame = app_box.text_frame
    app_frame.word_wrap = True
    
    p = app_frame.paragraphs[0]
    p.text = "🌐 APPLICATIONS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(10)
    
    applications = [
        ("🏫 Education", "Classroom attendance"),
        ("🏢 Corporate", "Office access control"),
        ("🔒 Security", "Facility monitoring"),
        ("🏦 Banking", "Fraud prevention"),
        ("📹 Surveillance", "Real-time monitoring")
    ]
    
    for app, desc in applications:
        p = app_frame.add_paragraph()
        p.text = app
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
        p = app_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.level = 1
        p.space_after = Pt(10)

def add_conclusion_slide():
    """Slide 7: Conclusion & Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, DARK_BLUE)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Conclusion & Call to Action"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Key Takeaways
    takeaway_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2))
    takeaway_frame = takeaway_box.text_frame
    takeaway_frame.word_wrap = True
    
    p = takeaway_frame.paragraphs[0]
    p.text = "✨ KEY TAKEAWAYS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(10)
    
    takeaways = [
        "🎯 Real-time face recognition with 95% accuracy",
        "⚡ GPU accelerated processing (150ms per frame)",
        "⚠️  Automatic email alerts for unknown persons",
        "🚀 Production-ready, deployed on GitHub"
    ]
    
    for takeaway in takeaways:
        p = takeaway_frame.add_paragraph()
        p.text = takeaway
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)
    
    # Demo invitation
    demo_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(1.2))
    demo_frame = demo_box.text_frame
    demo_frame.word_wrap = True
    
    p = demo_frame.paragraphs[0]
    p.text = "👉 LIVE DEMO AVAILABLE"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)
    
    p = demo_frame.add_paragraph()
    p.text = "Watch the system detect faces in real-time and trigger email alerts"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Footer - Contact
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.5))
    footer_frame = footer_box.text_frame
    footer_frame.word_wrap = True
    
    p = footer_frame.paragraphs[0]
    p.text = "Team: Saravanan, Praveen, Harish Raj, Tharun Kumar, Sathivel"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    
    p = footer_frame.add_paragraph()
    p.text = "GitHub: github.com/thirumalain186-collab/FaceID-Attendence-and-Alart-"
    p.font.size = Pt(12)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    
    p = footer_frame.add_paragraph()
    p.text = "🌟 Where AI Meets Automation 🌟"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

# Create all slides
print("Creating Slide 1: Title Slide...")
add_title_slide()

print("Creating Slide 2: Problem & Solution...")
add_problem_solution_slide()

print("Creating Slide 3: Architecture & Tech Stack...")
add_architecture_slide()

print("Creating Slide 4: Features & Performance...")
add_features_slide()

print("Creating Slide 5: Live Demo & Results...")
add_demo_slide()

print("Creating Slide 6: Innovation & Applications...")
add_innovation_slide()

print("Creating Slide 7: Conclusion & Call to Action...")
add_conclusion_slide()

# Save presentation
output_path = "Face_Recognition_Attendance_System_Presentation.pptx"
prs.save(output_path)
print(f"\nPresentation created successfully: {output_path}")
